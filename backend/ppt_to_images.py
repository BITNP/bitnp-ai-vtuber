#!/usr/bin/env python3
"""
PPT转图片工具
使用pywin32调用PowerPoint（Windows）进行转换，fallback到python-pptx
目前仅成功测试windows使用pywin32调用PowerPoint进行转换，未测试其他操作系统是否支持。

使用前需先安装依赖：
``` shell
cd backend
pip install pywin32 python-pptx matplotlib
```

运行以下命令进行转换：
``` shell
cd backend
python ppt_to_images.py <PPT文件路径> -o ../frontend/public/documents/slides
```
支持批量转换目录中的PPT文件：
``` shell
cd backend
python ppt_to_images.py <目录路径> -o ../frontend/public/documents/slides
```

转换完成后的ppt文件会在../frontend/public/documents/slides生成对应的图片文件。
"""
import os
import sys
import argparse
import glob
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt


def convert_pptx_to_images_powerpoint(pptx_path, output_folder="D:\\系统默认\\桌面\\bitnp-ai-vtuber\\frontend\\public\\documents\\slides"):
    """
    使用PowerPoint（Windows）将PPTX转换为高质量图片
    
    Args:
        pptx_path (str): PPTX文件路径
        output_folder (str): 输出图片目录
    """
    # 创建输出文件夹
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # 检查是否为Windows系统
    if not sys.platform.startswith('win'):
        print("PowerPoint转换仅支持Windows系统")
        return False
    
    try:
        import win32com.client
    except ImportError:
        print("未安装pywin32库，请先安装: pip install pywin32")
        return False
    
    print(f"正在使用PowerPoint转换PPT: {pptx_path}")
    
    try:
        # 创建PowerPoint应用对象
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True  # 设置为True确保能正确渲染
        
        # 转换为绝对路径
        pptx_path = os.path.abspath(pptx_path)
        output_folder = os.path.abspath(output_folder)
        
        # 打开PPT文件
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        
        # 设置导出格式为PNG
        format_type = 18  # 18 = PNG格式
        
        # 导出所有幻灯片
        presentation.Export(output_folder, "PNG", 1920, 1080)  # 1920x1080分辨率
        
        # 关闭PPT和PowerPoint
        presentation.Close()
        powerpoint.Quit()
        
        # 重命名文件（PowerPoint生成的命名格式为 "Slide1.png", "Slide2.png" 等）
        for filename in os.listdir(output_folder):
            if filename.startswith("Slide") and filename.endswith(".png"):
                # 提取页码
                try:
                    page_num = int(filename[5:-4])  # Slide1.png -> 1
                    new_filename = f"幻灯片{page_num}.PNG"  # 幻灯片1.PNG
                    old_path = os.path.join(output_folder, filename)
                    new_path = os.path.join(output_folder, new_filename)
                    
                    # 避免重复重命名
                    if old_path != new_path:
                        os.rename(old_path, new_path)
                        print(f"重命名: {filename} -> {new_filename}")
                except ValueError:
                    continue
        
        print("转换完成！")
        print(f"输出目录: {output_folder}")
        return True
        
    except Exception as e:
        print(f"PowerPoint转换失败: {e}")
        return False


def convert_pptx_to_images_python_pptx(pptx_path, output_folder="D:\\系统默认\\桌面\\bitnp-ai-vtuber\\frontend\\public\\documents\\slides"):
    """
    使用python-pptx库将PPTX转换为图片
    
    Args:
        pptx_path (str): PPTX文件路径
        output_folder (str): 输出图片目录
    """
    # 创建输出文件夹
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # 加载PPT
    prs = Presentation(pptx_path)
    
    print(f"正在使用python-pptx转换PPT: {pptx_path}")
    print(f"共 {len(prs.slides)} 页")
    
    for i, slide in enumerate(prs.slides):
        try:
            # 使用matplotlib保存（简单但可能不完美）
            fig = plt.figure(figsize=(16, 9))  # 16:9比例
            ax = fig.add_subplot(111)
            ax.set_xlim([0, 1])
            ax.set_ylim([0, 1])
            ax.axis('off')
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text:
                        # 获取文本位置和大小
                        left = shape.left / Inches(1)
                        top = shape.top / Inches(1)
                        width = shape.width / Inches(1)
                        height = shape.height / Inches(1)
                        
                        # 添加文本
                        ax.text(
                            left/16,  # 归一化坐标
                            1 - (top + height)/9,  # Y轴翻转
                            shape.text,
                            fontsize=12,
                            verticalalignment='top',
                            wrap=True
                        )
            
            # 保存图片
            output_path = os.path.join(output_folder, f"slide_{i+1:03d}.png")
            plt.savefig(output_path, dpi=150, bbox_inches='tight', pad_inches=0)
            plt.close(fig)
            
            print(f"已保存第 {i+1} 页: {output_path}")
            
        except Exception as e:
            print(f"第 {i+1} 页转换失败: {e}")
    
    print("转换完成！")
    print(f"输出目录: {output_folder}")
    return True


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPT转图片工具")
    parser.add_argument("path", help="PPTX文件路径或包含PPTX文件的目录路径")
    parser.add_argument("-o", "--output", default="D:\\系统默认\\桌面\\bitnp-ai-vtuber\\frontend\\public\\documents\\slides", help="输出图片目录")
    parser.add_argument("-r", "--recursive", action="store_true", help="递归查找子目录中的PPTX文件")
    
    args = parser.parse_args()
    
    # 确定要转换的文件列表
    pptx_files = []
    
    if os.path.isdir(args.path):
        # 如果是目录，查找所有PPTX文件
        pattern = "**/*.pptx" if args.recursive else "*.pptx"
        pptx_files = glob.glob(os.path.join(args.path, pattern), recursive=args.recursive)
        
        print(f"在目录 {args.path} 中找到 {len(pptx_files)} 个PPTX文件")
    else:
        # 如果是文件，直接添加到列表
        if args.path.lower().endswith('.pptx'):
            pptx_files = [args.path]
        else:
            print(f"错误：{args.path} 不是PPTX文件")
            sys.exit(1)
    
    # 转换所有找到的PPTX文件
    for pptx_file in pptx_files:
        print(f"\n=== 正在处理: {pptx_file} ===")
        
        # 创建输出目录
        output_dir = args.output
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # 使用PowerPoint进行转换
        result = convert_pptx_to_images_powerpoint(pptx_file, output_dir)
        
        if not result:
            # PowerPoint转换失败，尝试使用python-pptx
            print("\nPowerPoint转换失败，尝试使用python-pptx...")
            result = convert_pptx_to_images_python_pptx(pptx_file, output_dir)
        
        if not result:
            print(f"\n转换失败: {pptx_file}")
        else:
            print(f"\n转换成功: {pptx_file}")
    
    print(f"\n=== 处理完成，共处理 {len(pptx_files)} 个PPTX文件 ===")