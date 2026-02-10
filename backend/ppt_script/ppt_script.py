import asyncio
import os
import sys
import base64
import json
import requests
from abc import ABC, abstractmethod
from typing import List, Dict
import functools

from fontTools.ttx import process
from pptx import Presentation #  pip install python-pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from PIL import Image  #  pip install Pillow

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

try:
    from backend.llm_api.glm import GlmBot  # type: ignore
    from backend.tokens import get_token  # type: ignore
except ImportError:
    from llm_api.glm import GlmBot
    from tokens import get_token

#1.视觉识图
class AbstractVisionModel(ABC):
    @abstractmethod
    async def describe_image(self, base64_img:str) -> str | None:
        pass

class GlmVisionModel(AbstractVisionModel):
    def __init__(self, api_key: str,):
        self.api_key = api_key
        self.url = "https://open.bigmodel.cn/api/paas/v4/chat/completions"
        self.model_name = "glm-4v-flash"

    async def describe_image(self, base64_img:str) -> str | None:
        headers = {
            'Content-Type':'application/json',
            'Authorization': f'Bearer {self.api_key}'
        }

        prompt_text = """
                请分析这张 PPT 图片。
                1. 如果它是**有意义的内容**（如图表、数据截图、产品照片、核心文字），请详细描述其关键信息。
                2. 如果它是**无关装饰**（如简单的 LOGO、页码、纯背景纹理、无意义的线条或图标），请直接回复 "IGNORE_IMAGE"。
                3. **描述风格**：请用客观、详细的语言描述，不要加“这张图展示了”这种废话，直接说内容。
                """

        data = {
            "model": self.model_name,
            "messages": [{ "role": "user","content": [
                        {"type":"image_url","image_url":{"url": base64_img}},
                        {"type":"text","text": prompt_text}
            ]}],
            "stream": False,
            "temperature": 0.1
        }

        loop = asyncio.get_event_loop()
        try:
            func = functools.partial(requests.post, self.url, headers=headers, json=data)
            resp = await loop.run_in_executor(None, func)

            if resp.status_code != 200:
                return f"报错{resp.status_code}"
            content = resp.json()['choices'][0]['message']['content']

            if "IGNORE_IMAGE" in content:
                return None
            return content

        except Exception as e:
            return f"报错{e}"

#2.ppt解析
class PPTparser:
    def __init__(self, ppt_path:str):
        self.ppt_path = ppt_path

    def get_iter_shapes(self,shapes):
         for shape in shapes:
             if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                 yield from self.get_iter_shapes(shape.shapes)
             else:
                 yield shape

    def process_image(self, image_blob):
        """转JPG+压缩+添加Header"""
        try:
            img = Image.open(io.BytesIO(image_blob))
            if img.mode != 'RGB': img = img.convert('RGB')
            img.thumbnail((1024,1024))

            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=85)
            b64 = base64.b64encode(buf.getvalue()).decode('utf-8')
            return f"data:image/jpeg;base64,{b64}"
        except:
            return None

    def extract_content(self) -> List[Dict]:
#        if not os.path.exists(self.ppt_path):
#            raise FileNotFoundError(f"PPT不存在")
        prs = Presentation(self.ppt_path)
        slides_data = []

        for i,slide in enumerate(prs.slides):
            text_runs = []
            images_bs64 = []

            for shape in self.get_iter_shapes(slide.shapes):
                if shape.has_text_frame and shape.text.strip():
                    text_runs.append(shape.text.strip())

                if hasattr(shape, "image"):
                    if len(shape.image.blob) > 3*1024:
                        processed = self.process_image(shape.image.blob)
                        if processed:
                            images_bs64.append(processed)

            page_text = "\n".join(text_runs).strip()
            if not page_text:
                page_text = "本页无文字，根据图片内容"

            slides_data.append({
                "page": i+1,
                "text": page_text,
                "images": images_bs64,
                })
        return slides_data


#3.核心处理

async def process_slide_vision(page_idx,images,bot,sem) -> str:
    if not images: return ""

    valid_descriptions = []
    async def recogmnnize(img):
        async with sem:
            return await bot.describe_image(img)

    results = await asyncio.gather(*[recogmnnize(img) for img in images])

    for result in results:
        if result:
            valid_descriptions.append(result)

    if not valid_descriptions: return ""

    return "\n".join([f"图片{i+1}:{d}"for i,d in enumerate(valid_descriptions)])

async  def generate_ppt_scripts(ppt_path:str, vision_config:Dict, text_gen_api_key:str):
    parser = PPTparser(ppt_path)
    slides = parser.extract_content()
    total_pages = len(slides)

    p = vision_config.get("provider","glm")
    if p == 'glm':
        vision_bot = GlmVisionModel(vision_config["api_key"])
    else:
        raise ValueError(f"Unsupported provider:")

    #1
    sem = asyncio.Semaphore(3)
    prepared_data = []
    vision_tasks = []

    for slide in slides:
        task = process_slide_vision(slide['page'], slide['images'], vision_bot, sem)
        vision_tasks.append(task)

        prepared_data.append({
            "page":slide['page'],
            "text":slide['text'],
            "vision":None
        })

    vision_results = await asyncio.gather(*vision_tasks)

    outline_lines = []
    for i,v_res in enumerate(vision_results):
        prepared_data[i]['vision'] = v_res
        preview = prepared_data[i]['text'].replace('\n', ' ')
        outline_lines.append(f"第{i+1}页: {preview}")
#        print(f"第{i+1}页，有效信息：{'have'if v_res and 'ignore' not in v_res else 'no'}")

    global_outline = "\n".join(outline_lines)

    #2

    system_prompt = f"""
        你现在是“树莓娘”，是网络开拓者协会（网协，北理工学生组织）的看板娘。
        你正在进行一场技术分享会的直播，主题是关于PPT内容的讲解。
        全局：{global_outline}

        *** 核心风格 ***：
        1. 不要机械地朗读PPT上的文字！要结合【视觉描述】和【文字内容】。   
        2. 用口语化的表达，不要读得太书面化。
        3. 表情动作：说话时必须自然地穿插表情指令，仅限使用：[点头] [摇头] [wink]。平均每段话使用1~2个。

        *** 强制格式规范（必须严格遵守） ***：
        1. PPT切换指令：每一页讲稿的**最开头**，必须严格输出 `[PPT_x]` （x是当前页码）。
           - 正确示例：`[PPT_1] 大家好！我是树莓娘...`
           - 错误示例：`好的，第1页：大家好...`
        2. 严禁 Markdown：绝对不要使用 `**加粗**`、`# 标题`、`- 列表` 等符号，直接用自然语言表达。
        3. 逻辑衔接：严禁使用“下一页是”、“接下来看下一页”、“好的”、"好呀"这种报幕词。要用内容逻辑自然过渡。
           - 避免在开头说当前页码。

        *** 视觉讲解逻辑 ***：
        - 如果视觉描述显示是“图表/代码/截图”，请具体描述图里的细节（例如“大家看这行代码...”）。
        - 如果视觉描述显示是“装饰/LOGO”，则忽略图片，专注于文字概念的讲解。
        """

    text_bot = GlmBot(
        token = text_gen_api_key,
        model_name = 'glm-4-flash',
        system_prompt = system_prompt
    )

    full_results = []
    last_script = "(开场)"

    for i, data in enumerate(prepared_data):
        page = data['page']

        next_info = "(这是最后一页，请进行总结并和大家说再见)"
        if i < total_pages - 1:
            next_data = prepared_data[i + 1]
            next_txt = next_data['text'].replace('\n', ' ')
            next_vision = next_data['vision'] if next_data['vision'] else "无重点图片"
            next_info = f"下一页内容预览：文字提到了“{next_txt}”，视觉包含“{next_vision}”"

        prompt = f"""
            当前任务：请生成第 {page} 页的讲解台词。

            【当前页输入信息】
            1. PPT文字内容：
            {data['text']}

            2. PPT图片描述：
            {data['vision']}

            3. 上一页讲稿：
            “{last_script if last_script else '无'}”

            4. {next_info}

            【生成要求】
            - 必须以 `[PPT_{page}]` 开头。
            - 语气活泼，禁止 Markdown 符号 和 emoji 符号。
            - 避免使用“下一页”、“接下来”等报幕词，避免在开头说当前页码。
            - 不要说“好呀”、“好的”等口头禅。
            - 避免在开头回顾上一页内容，要直接切入主题；上一页讲稿仅作为Callback参考。
            - 拒绝平铺直叙：不要把PPT上的字都念一遍！挑一个重点深入讲。
            - 包含 1-2 个表情（[点头]/[摇头]/[wink]）。
            - 如果当前页文字极少（如仅有标题），请结合视觉描述或主题进行发挥，不要只说一句话。
            """

        text_bot.messages = []
        text_bot.append_context(prompt)

        script = await text_bot.respond_to_context()
        last_script = script

        full_results.append({
            "page": page,
            "text": data['text'],
            "vision": data['vision'],
            "script": script
        })

    save_files(ppt_path, full_results)

def save_files(ppt_path, results):
    """
    保存四个文件:
    1.*_json: 每页的PPT文字、视觉描述和讲解台词
    2.*_ppt_text: 每页的文字内容
    3.*_scripts: 每页的讲解
    4.*_vision: 每页图片的视觉描述
    保存在PPT文件所在目录下自动创建的名为 generated_scripts 的文件夹，
    """
    out_dir = os.path.join(os.path.dirname(os.path.abspath(ppt_path)), 'generated_scripts')
    if not os.path.exists(out_dir):
         os.makedirs(out_dir)

    base = os.path.splitext(os.path.basename(ppt_path))[0]

    json_path = os.path.join(out_dir, f"{base}_data.json")
    with open(json_path, 'w', encoding='utf-8') as f:
          json.dump(results, f, ensure_ascii=False, indent=4)

    script_txt_path = os.path.join(out_dir, f"{base}_scripts.txt")
    with open(script_txt_path, 'w', encoding='utf-8') as f:
        for item in results:
            # f.write(f"第{item['page']}页\n")
            f.write(f"{item['script']}\n\n")
            # f.write("\n\n")

    vision_txt_path = os.path.join(out_dir, f"{base}_vision.txt")
    with open(vision_txt_path, 'w', encoding='utf-8') as f:
        for item in results:
            f.write(f" 第 {item['page']} 页 \n")
            f.write(f"{item['vision'] if item['vision'] else '(无有效视觉信息)'}\n")
            f.write("\n")

    ppt_text_path = os.path.join(out_dir, f"{base}_ppt_text.txt")
    with open(ppt_text_path, 'w', encoding='utf-8') as f:
        for item in results:
            f.write(f" 第 {item['page']} 页 \n")
            f.write(f"{item['text']}\n")
            f.write("\n")

from argparse import ArgumentParser

def args_parse():
    parser = ArgumentParser(description="PPT讲稿生成器")
    parser.add_argument("ppt_path", type=str, help="PPT文件路径")
    return parser.parse_args()

if __name__ == "__main__":
    from tokens import get_token
    
    VISION_CFG = {
        "provider": "glm",
        "api_key": get_token('glm')
    }
    TEXT_KEY = get_token('glm')
    args = args_parse()
    PPT_FILE = args.ppt_path
    if os.path.exists(PPT_FILE):
        asyncio.run(generate_ppt_scripts(PPT_FILE, VISION_CFG, TEXT_KEY))
    else:
        print("文件不存在")
